import os
from pptx import Presentation
import fitz  # PyMuPDF，用于 PDF 解析
from PIL import Image
import pytesseract

# 显式指定 Tesseract 的安装路径（强烈推荐）
pytesseract.pytesseract.tesseract_cmd = r"D:\WorkSoftware\Tesseract-OCR\tesseract.exe"

def extract_text_from_file(file_path, file_type):
    try:
        if file_type == 'pdf':
            return extract_from_pdf(file_path)
        elif file_type == 'pptx':
            return extract_from_pptx(file_path)
        else:
            return '[暂不支持该类型提取]'
    except Exception as e:
        return f'[提取失败: {str(e)}]'

# 文本型 PDF + 图片型 PDF（自动混合处理）
def extract_from_pdf(path):
    text = ''
    doc = fitz.open(path)
    for i, page in enumerate(doc):
        # 先尝试提取可选文字
        page_text = page.get_text()
        text += page_text

        # 如果该页无文本，则尝试截图 + OCR
        if not page_text.strip():
            # 将 PDF 页转为图像
            pix = page.get_pixmap(dpi=500)
            image_path = f"temp_page_{i}.png"
            pix.save(image_path)

            # OCR 识别图像文字
            image = Image.open(image_path)
            ocr_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            text += ocr_text + '\n'

            # 清理临时图像
            os.remove(image_path)
    return text.strip()

def extract_from_pptx(path):
    prs = Presentation(path)
    text = ''

    # 提取文字内容
    for slide in prs.slides:
        for shape in slide.shapes:
            # 文字框
            if hasattr(shape, "text"):
                text += shape.text + '\n'

            # 图片 OCR
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                image_bytes = image.blob
                image_stream = io.BytesIO(image_bytes)

                # 保存为临时文件用于 OCR
                tmp_image_path = 'temp_slide_image.png'
                with open(tmp_image_path, 'wb') as f:
                    f.write(image_stream.read())

                try:
                    img = Image.open(tmp_image_path)
                    ocr_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    text += '\n[图片识别文字]:\n' + ocr_text + '\n'
                except Exception as e:
                    text += f'\n[图片OCR失败: {str(e)}]\n'

                # 删除临时图片
                os.remove(tmp_image_path)

    return text.strip()
